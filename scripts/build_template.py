# -*- coding: utf-8 -*-
"""스타일 베이스(assets/ppt-template.pptx)를 스크립트로 재생성(복각).

originals/ppt-template-old.pptx 에서 마스터/테마/레이아웃/배경/로고(your-logo)만 가져오고,
본문 콘텐츠 슬라이드는 모두 제거한 뒤, 쇼케이스용 표지(변주 A) + 감사 페이지만 얹는다.
skill의 render.build()는 이 파일을 로드해 슬라이드를 비우고 새 덱을 생성하므로,
여기 표지/감사는 'ppt-template.pptx를 직접 열었을 때 스타일을 보여주는 견본'이다.
"""
import sys, os, json
sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))
import render
from pptx import Presentation

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
OLD = os.path.join(ROOT, "originals", "ppt-template-old.pptx")
OUT = os.path.join(ROOT, "assets", "ppt-template.pptx")
CFG = os.path.join(ROOT, "config.json")

accent = "1456F0"
if os.path.exists(CFG):
    accent = json.load(open(CFG, encoding="utf-8")).get("accent", accent)
render.configure(accent)

prs = Presentation(OLD)
render._delete_all_slides(prs)
layout = prs.slide_layouts[0]

cover = prs.slides.add_slide(layout)
render.cover(cover, {
    "eyebrow": "YOUR REPORT",
    "title1": "프레젠테이션 제목을",
    "title2": "여기에 입력하세요",
    "subtitle": "부제 · 한 줄 요약을 입력하는 영역입니다",
    "foot": "2026 · Your Organization",
})

thx = prs.slides.add_slide(layout)
render.thanks(thx, {
    "eyebrow": "THANK YOU",
    "title": "감사합니다",
    "subtitle": "끝까지 읽어주셔서 감사합니다",
    "foot": "contact@your-logo.com   ·   your-logo",
})

os.makedirs(os.path.dirname(OUT), exist_ok=True)
prs.save(OUT)
print("재생성 완료:", OUT, "| accent #%s" % accent)
