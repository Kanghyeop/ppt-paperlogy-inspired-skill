# -*- coding: utf-8 -*-
"""ppt-skill 통일 렌더러 (python-pptx). 슬라이드 JSON을 단일 디자인 시스템으로 그린다.
디자인은 고정(회색 배경·Pretendard·프레임), 포인트 컬러(accent)만 config로 교체 가능."""
import re, math
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import CategoryChartData
from pptx.oxml.ns import qn

# ---------- 색/지오메트리 ----------
BLUE, BLUE2, BLUE3, BLUEL = "1456F0", "60A5FA", "17437D", "BFDBFE"
INK, SLATE, GRAY, LINE = "222222", "45515E", "8E8E93", "E5E7EB"
WHITE, CARDG = "FFFFFF", "F2F3F5"
SW, SH = 12192000, 6858000
MX = 457200
BODY_X, BODY_Y, BODY_W = MX, 2255520, 11277600
BODY_Y2 = 6080000
BODY_H = BODY_Y2 - BODY_Y
HANGUL = re.compile('[가-힣]')

def C(s): return RGBColor.from_string(s)

def configure(accent=None):
    """포인트 컬러만 교체(나머지 디자인은 고정). accent: 'RRGGBB' 또는 '#RRGGBB'."""
    global BLUE
    if accent:
        BLUE = accent.lstrip("#").upper()

# ---------- 폰트 자동 맞춤 ----------
def _text_units(s):
    """글자 폭 추정치(전각=1.0, 반각=0.55)."""
    u = 0.0
    for ch in s:
        u += 1.0 if ord(ch) > 0x1100 else 0.55
    return u

def fit_font(text, box_w_emu, max_lines, base_pt, min_pt):
    """가장 긴 단어가 한 줄에 들어가고 전체가 max_lines 이내가 되는 최대 폰트(pt)."""
    text = text or ""
    if not text.strip():
        return base_pt
    total = _text_units(text)
    longest = max((_text_units(t) for t in text.split()), default=total)
    for size in range(int(base_pt), int(min_pt) - 1, -1):
        cap = box_w_emu / (size * 12700 * 1.08)   # 한 줄 용량(units), 단어 줄바꿈 여유 반영
        if cap <= 0:
            continue
        if longest <= cap and math.ceil(total / cap) <= max_lines:
            return size
    return int(min_pt)

def _set_shrink(tf):
    """넘치면 자동 축소(normAutofit) - 안전망."""
    try:
        bodyPr = tf._txBody.find(qn('a:bodyPr'))
        if bodyPr is None:
            return
        for tag in ('a:normAutofit', 'a:spAutoFit', 'a:noAutofit'):
            e = bodyPr.find(qn(tag))
            if e is not None:
                bodyPr.remove(e)
        bodyPr.append(bodyPr.makeelement(qn('a:normAutofit'), {}))
    except Exception:
        pass

# ---------- 텍스트 ----------
def _style_run(run, size, bold, color):
    run.font.size = Pt(size)
    run.font.bold = bool(bold)
    run.font.name = "Pretendard"
    run.font.color.rgb = C(color)
    rPr = run._r.get_or_add_rPr()
    spc = round(size * -5) if (run.text and HANGUL.search(run.text)) else 0
    rPr.set('spc', str(spc))

def fill_tf(tf, paras, anchor=MSO_ANCHOR.TOP, wrap=True, shrink=False):
    """para = (text,size,bold,color,align) 또는 (text,size,bold,color,align,space_before_pt)."""
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ('margin_left','margin_right','margin_top','margin_bottom'):
        setattr(tf, m, 0)
    first = True
    for p in paras:
        text, size, bold, color, align = p[0], p[1], p[2], p[3], p[4]
        sb = p[5] if len(p) > 5 else None
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.alignment = align
        if sb is not None:
            para.space_before = Pt(sb)
            para.space_after = Pt(0)
        run = para.add_run()
        run.text = text
        _style_run(run, size, bold, color)
    if shrink:
        _set_shrink(tf)
    return tf

def textbox(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP, wrap=True, shrink=True):
    tb = slide.shapes.add_textbox(Emu(x), Emu(y), Emu(w), Emu(h))
    fill_tf(tb.text_frame, paras, anchor, wrap, shrink)
    return tb

def box(slide, x, y, w, h, fill=None, line=None, rounded=False, radius=0.055):
    shp = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE,
        Emu(x), Emu(y), Emu(w), Emu(h))
    if fill:
        shp.fill.solid(); shp.fill.fore_color.rgb = C(fill)
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = C(line); shp.line.width = Pt(1)
    else:
        shp.line.fill.background()
    try: shp.shadow.inherit = False
    except Exception: pass
    if rounded:
        try: shp.adjustments[0] = radius
        except Exception: pass
    return shp

L, Cn, R = PP_ALIGN.LEFT, PP_ALIGN.CENTER, PP_ALIGN.RIGHT

# ---------- 프레임 (헤더/제목/부제/페이지/캡션) ----------
def frame(slide, header, page_no, total, caption, title, subtitle):
    textbox(slide, MX, 365760, 6400800, 274320,
            [(header, 11, True, GRAY, L)], anchor=MSO_ANCHOR.MIDDLE)
    textbox(slide, MX, 914400, 11277295, 685800,
            [(title, 30, True, INK, L)], anchor=MSO_ANCHOR.TOP)
    if subtitle:
        textbox(slide, MX, 1490472, 11277295, 365760,
                [(subtitle, 14, False, SLATE, L)], anchor=MSO_ANCHOR.TOP)
    textbox(slide, MX, 6446520, 2743200, 228600,
            [("%02d / %02d" % (page_no, total), 9, False, GRAY, L)], anchor=MSO_ANCHOR.MIDDLE)
    if caption:
        textbox(slide, 3505200, 6446520, 8229600, 228600,
                [(caption, 9, False, GRAY, R)], anchor=MSO_ANCHOR.MIDDLE)

# ---------- 차트 공통 ----------
def _style_chart(chart, legend=False, legend_pos='b'):
    chart.has_title = False
    chart.font.name = "Pretendard"
    chart.font.size = Pt(10)
    chart.font.color.rgb = C(SLATE)
    chart.has_legend = legend
    if legend:
        chart.legend.position = {'b': XL_LEGEND_POSITION.BOTTOM, 'r': XL_LEGEND_POSITION.RIGHT}[legend_pos]
        chart.legend.include_in_layout = False
        chart.legend.font.size = Pt(10); chart.legend.font.color.rgb = C(SLATE)

def _axis_clean(chart, vis_value_grid=True):
    ca = chart.category_axis
    ca.has_major_gridlines = False
    ca.tick_labels.font.size = Pt(10); ca.tick_labels.font.color.rgb = C(SLATE)
    ca.format.line.color.rgb = C(LINE)
    try:
        va = chart.value_axis
        va.has_major_gridlines = vis_value_grid
        if vis_value_grid:
            va.major_gridlines.format.line.color.rgb = C(LINE)
            va.major_gridlines.format.line.width = Pt(0.5)
        va.tick_labels.font.size = Pt(9); va.tick_labels.font.color.rgb = C(GRAY)
        va.format.line.fill.background()
    except Exception:
        pass

def _chart_frame(slide):
    return (BODY_X, BODY_Y, BODY_W, BODY_H)

# ---------- 블록 렌더러 ----------
def b_stat_cards(slide, blk):
    cards = blk["cards"][:4]
    n = len(cards)
    gap = 274320
    cw = (BODY_W - gap * (n - 1)) // n
    ch = 2360000
    cy = BODY_Y + (BODY_H - ch) // 2
    pad = 300000
    iw = cw - 2 * pad
    # 카드 간 value 폰트 통일: 가장 긴 value 기준 최소값
    vsize = min(fit_font(str(c.get("value", "")), iw, 2, 40, 18) for c in cards)
    for i, c in enumerate(cards):
        cx = BODY_X + i * (cw + gap)
        box(slide, cx, cy, cw, ch, fill=WHITE, rounded=True)
        val = str(c.get("value", "")); label = str(c.get("label", "")); sub = str(c.get("sub", ""))
        paras = [(val, vsize, True, BLUE, L)]
        if label:
            paras.append((label, 14, True, INK, L, 10))   # 앞 여백 10pt
        if sub:
            paras.append((sub, 10, False, GRAY, L, 4))
        # 한 프레임에 흘려서(value→label→sub) 절대 겹치지 않게, 세로 중앙 정렬
        textbox(slide, cx + pad, cy + 200000, iw, ch - 400000, paras,
                anchor=MSO_ANCHOR.MIDDLE)

def b_numbered(slide, blk):
    items = blk["items"][:5]
    n = len(items)
    gap = 137160
    rh = (BODY_H - gap * (n - 1)) // n
    for i, it in enumerate(items):
        y = BODY_Y + i * (rh + gap)
        textbox(slide, BODY_X, y, 900000, rh,
                [("%02d" % (i + 1), 22, True, BLUE, L)], anchor=MSO_ANCHOR.MIDDLE)
        box(slide, BODY_X + 950000, y + rh // 2 - 5000, 6000, rh - 200000, fill=LINE)
        textbox(slide, BODY_X + 1120000, y, 3400000, rh,
                [(str(it.get("title", "")), 15, True, INK, L)], anchor=MSO_ANCHOR.MIDDLE)
        textbox(slide, BODY_X + 4700000, y, BODY_W - 4700000, rh,
                [(str(it.get("desc", "")), 12, False, SLATE, L)], anchor=MSO_ANCHOR.MIDDLE)

def b_timeline(slide, blk):
    steps = blk["steps"][:6]
    n = len(steps)
    axis_y = BODY_Y + BODY_H // 2
    box(slide, BODY_X, axis_y - 3000, BODY_W, 6000, fill=LINE)
    seg = BODY_W // n
    for i, s in enumerate(steps):
        cx = BODY_X + seg * i + seg // 2
        dot = 130000
        box(slide, cx - dot // 2, axis_y - dot // 2, dot, dot, fill=BLUE, rounded=True, radius=0.5)
        # 위: label
        textbox(slide, cx - seg // 2 + 40000, axis_y - 1320000, seg - 80000, 320000,
                [(str(s.get("label", "")), 13, True, BLUE, Cn)], anchor=MSO_ANCHOR.BOTTOM)
        textbox(slide, cx - seg // 2 + 40000, axis_y - 980000, seg - 80000, 560000,
                [(str(s.get("title", "")), 12, True, INK, Cn)], anchor=MSO_ANCHOR.BOTTOM)
        # 아래: desc
        textbox(slide, cx - seg // 2 + 40000, axis_y + 230000, seg - 80000, 900000,
                [(str(s.get("desc", "")), 10, False, SLATE, Cn)], anchor=MSO_ANCHOR.TOP)

def b_process(slide, blk):
    steps = blk["steps"][:5]
    n = len(steps)
    arrow = 200000
    gap = 120000
    bw = (BODY_W - (n - 1) * (arrow + gap)) // n
    bh = 2100000
    y = BODY_Y + (BODY_H - bh) // 2
    for i, s in enumerate(steps):
        x = BODY_X + i * (bw + arrow + gap)
        box(slide, x, y, bw, bh, fill=WHITE, rounded=True)
        textbox(slide, x + 200000, y + 230000, bw - 400000, 360000,
                [("STEP %d" % (i + 1), 9, True, BLUE, L)])
        textbox(slide, x + 200000, y + 620000, bw - 400000, 560000,
                [(str(s.get("title", "")), 13, True, INK, L)])
        textbox(slide, x + 200000, y + 1180000, bw - 400000, 760000,
                [(str(s.get("desc", "")), 10, False, SLATE, L)])
        if i < n - 1:
            ax = x + bw + gap // 2
            tri = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, Emu(ax), Emu(y + bh // 2 - 90000), Emu(arrow), Emu(180000))
            tri.fill.solid(); tri.fill.fore_color.rgb = C(BLUEL)
            tri.line.fill.background()
            try: tri.shadow.inherit = False
            except Exception: pass

def b_comparison(slide, blk):
    gap = 320040
    pw = (BODY_W - gap) // 2
    for i, side in enumerate((blk["left"], blk["right"])):
        x = BODY_X + i * (pw + gap)
        accent = BLUE if i == 1 else SLATE
        box(slide, x, BODY_Y, pw, BODY_H, fill=WHITE, rounded=True)
        box(slide, x, BODY_Y, pw, 150000, fill=accent, rounded=False)
        textbox(slide, x + 300000, BODY_Y + 320000, pw - 600000, 480000,
                [(str(side.get("title", "")), 17, True, INK, L)])
        pts = side.get("points", [])[:4]
        py = BODY_Y + 980000
        for p in pts:
            box(slide, x + 320000, py + 70000, 90000, 90000, fill=accent, rounded=True, radius=0.5)
            textbox(slide, x + 520000, py, pw - 820000, 560000,
                    [(str(p), 12, False, SLATE, L)])
            py += 640000

def b_table(slide, blk):
    headers = blk["headers"]
    rows = blk["rows"]
    nr, nc = len(rows) + 1, len(headers)
    th = min(BODY_H, 520000 + (len(rows)) * 560000)
    gt = slide.shapes.add_table(nr, nc, Emu(BODY_X), Emu(BODY_Y), Emu(BODY_W), Emu(th)).table
    gt.first_row = False; gt.horz_banding = False
    try:
        gt.first_col = False; gt.last_row = False; gt.last_col = False; gt.vert_banding = False
    except Exception: pass
    for j, htxt in enumerate(headers):
        cell = gt.cell(0, j)
        cell.fill.solid(); cell.fill.fore_color.rgb = C(INK)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = Emu(110000); cell.margin_right = Emu(110000)
        fill_tf(cell.text_frame, [(str(htxt), 11, True, WHITE, L)])
    for i, row in enumerate(rows):
        for j in range(nc):
            cell = gt.cell(i + 1, j)
            cell.fill.solid(); cell.fill.fore_color.rgb = C(WHITE if i % 2 == 0 else CARDG)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = Emu(110000); cell.margin_right = Emu(110000)
            val = str(row[j]) if j < len(row) else ""
            fill_tf(cell.text_frame, [(val, 11, j == 0, INK if j == 0 else SLATE, L)])

def b_bullets(slide, blk):
    cols = blk["columns"][:2]
    gap = 360000
    cw = (BODY_W - gap * (len(cols) - 1)) // len(cols)
    for ci, col in enumerate(cols):
        x = BODY_X + ci * (cw + gap)
        y = BODY_Y
        for it in col[:3]:
            box(slide, x, y + 60000, 110000, 110000, fill=BLUE, rounded=True, radius=0.3)
            textbox(slide, x + 260000, y, cw - 260000, 420000,
                    [(str(it.get("head", "")), 14, True, INK, L)])
            textbox(slide, x + 260000, y + 430000, cw - 260000, 760000,
                    [(str(it.get("body", "")), 11, False, SLATE, L)])
            y += 1200000

def b_kpi_progress(slide, blk):
    items = blk["items"][:5]
    n = len(items)
    gap = 150000
    rh = (BODY_H - gap * (n - 1)) // n
    track_x = BODY_X + 3000000
    track_w = BODY_W - 3000000 - 900000
    for i, it in enumerate(items):
        y = BODY_Y + i * (rh + gap)
        mid = y + rh // 2
        textbox(slide, BODY_X, y, 2900000, rh,
                [(str(it.get("label", "")), 13, True, INK, L)], anchor=MSO_ANCHOR.MIDDLE)
        bar_h = 150000
        box(slide, track_x, mid - bar_h // 2, track_w, bar_h, fill=LINE, rounded=True, radius=0.5)
        pct = max(0, min(100, float(it.get("pct", 0))))
        fw = int(track_w * pct / 100)
        if fw > 0:
            box(slide, track_x, mid - bar_h // 2, fw, bar_h, fill=BLUE, rounded=True, radius=0.5)
        textbox(slide, track_x + track_w + 60000, y, 840000, rh,
                [("%d%%" % round(pct), 13, True, BLUE, L)], anchor=MSO_ANCHOR.MIDDLE)
        if it.get("note"):
            textbox(slide, track_x, mid + 80000, track_w, 300000,
                    [(str(it["note"]), 9, False, GRAY, L)])

def b_callout(slide, blk):
    ph = 1700000
    box(slide, BODY_X, BODY_Y, BODY_W, ph, fill=BLUE, rounded=True)
    big = str(blk.get("big", ""))
    bsize = fit_font(big, BODY_W - 840000, 2, 26, 16)
    textbox(slide, BODY_X + 420000, BODY_Y, BODY_W - 840000, ph,
            [(big, bsize, True, WHITE, L)], anchor=MSO_ANCHOR.MIDDLE)
    pts = blk.get("points", [])[:3]
    n = max(1, len(pts))
    gap = 300000
    cw = (BODY_W - gap * (n - 1)) // n
    cy = BODY_Y + ph + 300000
    chh = BODY_Y2 - cy
    for i, p in enumerate(pts):
        x = BODY_X + i * (cw + gap)
        box(slide, x, cy, cw, chh, fill=WHITE, rounded=True)
        textbox(slide, x + 240000, cy + 200000, cw - 480000, 360000,
                [("0%d" % (i + 1), 13, True, BLUE, L)])
        textbox(slide, x + 240000, cy + 580000, cw - 480000, chh - 700000,
                [(str(p), 12, False, SLATE, L)])

# ----- 차트 -----
def _palette(n):
    base = [BLUE, BLUE2, BLUE3, BLUEL, GRAY, SLATE]
    return [base[i % len(base)] for i in range(n)]

def b_column_chart(slide, blk):
    cd = CategoryChartData()
    cd.categories = blk["categories"]
    for s in blk["series"]:
        cd.add_series(s["name"], tuple(s["values"]))
    x, y, w, h = _chart_frame(slide)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Emu(x), Emu(y), Emu(w), Emu(h), cd)
    chart = gf.chart
    _style_chart(chart, legend=len(blk["series"]) > 1)
    _axis_clean(chart, vis_value_grid=True)
    for i, plot_series in enumerate(chart.series):
        plot_series.format.fill.solid(); plot_series.format.fill.fore_color.rgb = C(_palette(len(chart.series))[i])
    if len(blk["series"]) == 1:
        plot = chart.plots[0]; plot.has_data_labels = True
        plot.data_labels.font.size = Pt(10); plot.data_labels.font.bold = True
        plot.data_labels.font.color.rgb = C(INK)
        plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        plot.gap_width = 80

def b_bar_chart(slide, blk):
    cd = CategoryChartData()
    cd.categories = blk["categories"]
    for s in blk["series"]:
        cd.add_series(s["name"], tuple(s["values"]))
    x, y, w, h = _chart_frame(slide)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, Emu(x), Emu(y), Emu(w), Emu(h), cd)
    chart = gf.chart
    _style_chart(chart, legend=len(blk["series"]) > 1)
    _axis_clean(chart, vis_value_grid=False)
    for i, ps in enumerate(chart.series):
        ps.format.fill.solid(); ps.format.fill.fore_color.rgb = C(_palette(len(chart.series))[i])
    if len(blk["series"]) == 1:
        plot = chart.plots[0]; plot.has_data_labels = True
        plot.data_labels.font.size = Pt(10); plot.data_labels.font.bold = True
        plot.data_labels.font.color.rgb = C(INK)
        plot.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
        plot.gap_width = 70

def b_line_chart(slide, blk):
    cd = CategoryChartData()
    cd.categories = blk["categories"]
    for s in blk["series"]:
        cd.add_series(s["name"], tuple(s["values"]))
    x, y, w, h = _chart_frame(slide)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.LINE_MARKERS, Emu(x), Emu(y), Emu(w), Emu(h), cd)
    chart = gf.chart
    _style_chart(chart, legend=len(blk["series"]) > 1)
    _axis_clean(chart, vis_value_grid=True)
    pal = _palette(len(chart.series))
    for i, ps in enumerate(chart.series):
        ps.format.line.color.rgb = C(pal[i]); ps.format.line.width = Pt(2.5)
        ps.smooth = False
        try:
            ps.marker.style = 8  # circle
            ps.marker.format.fill.solid(); ps.marker.format.fill.fore_color.rgb = C(pal[i])
            ps.marker.format.line.color.rgb = C(WHITE)
        except Exception: pass

def b_donut_chart(slide, blk):
    cd = CategoryChartData()
    cd.categories = blk["labels"]
    cd.add_series("share", tuple(blk["values"]))
    x, y, w, h = _chart_frame(slide)
    gf = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, Emu(x), Emu(y), Emu(w), Emu(h), cd)
    chart = gf.chart
    _style_chart(chart, legend=True, legend_pos='r')
    pal = _palette(len(blk["values"]))
    series = chart.series[0]
    for idx, point in enumerate(series.points):
        point.format.fill.solid(); point.format.fill.fore_color.rgb = C(pal[idx])
        point.format.line.color.rgb = C(WHITE); point.format.line.width = Pt(2)
    plot = chart.plots[0]
    plot.has_data_labels = True
    plot.data_labels.number_format = '0"%"'; plot.data_labels.number_format_is_linked = False
    plot.data_labels.font.size = Pt(10); plot.data_labels.font.bold = True
    plot.data_labels.font.color.rgb = C(WHITE)

def b_image(slide, blk):
    path = blk.get("path", "")
    cap = str(blk.get("caption", ""))
    area_h = BODY_H - (340000 if cap else 0)
    try:
        from PIL import Image
        iw, ih = Image.open(path).size
        ar = iw / ih if ih else 16 / 9
    except Exception:
        ar = 16 / 9
    w = BODY_W; h = int(w / ar)
    if h > area_h:
        h = area_h; w = int(h * ar)
    x = BODY_X + (BODY_W - w) // 2
    y = BODY_Y + (area_h - h) // 2
    try:
        slide.shapes.add_picture(path, Emu(x), Emu(y), Emu(w), Emu(h))
    except Exception:
        box(slide, BODY_X, BODY_Y, BODY_W, area_h, fill=CARDG, rounded=True)
        textbox(slide, BODY_X, BODY_Y, BODY_W, area_h,
                [("[이미지를 찾을 수 없음: %s]" % path, 12, False, GRAY, Cn)], anchor=MSO_ANCHOR.MIDDLE)
    if cap:
        textbox(slide, BODY_X, BODY_Y + area_h + 90000, BODY_W, 280000,
                [(cap, 10, False, GRAY, Cn)], anchor=MSO_ANCHOR.TOP)

BLOCKS = {
    "stat_cards": b_stat_cards, "numbered": b_numbered, "timeline": b_timeline,
    "process": b_process, "comparison": b_comparison, "table": b_table,
    "bullets": b_bullets, "kpi_progress": b_kpi_progress, "callout": b_callout,
    "column_chart": b_column_chart, "bar_chart": b_bar_chart,
    "line_chart": b_line_chart, "donut_chart": b_donut_chart, "image": b_image,
}

# ---------- 표지 (변주 A) ----------
def cover(slide, d):
    t1 = str(d.get("title1", "")); t2 = str(d.get("title2", ""))
    lines = [t for t in (t1, t2) if t]
    tw = 10500000
    tsize = min([fit_font(t, tw, 1, 54, 30) for t in lines] or [54])
    textbox(slide, MX, 1828800, 8000000, 274320,
            [(str(d.get("eyebrow", "")), 11, True, BLUE, L)], anchor=MSO_ANCHOR.MIDDLE)
    textbox(slide, MX, 2240280, tw, 1828800,
            [(t, tsize, True, INK, L) for t in lines])
    box(slide, MX, 4343400, 640080, 64008, fill=BLUE)
    textbox(slide, MX, 4617720, 9500000, 365760,
            [(str(d.get("subtitle", "")), 18, False, SLATE, L)])
    textbox(slide, MX, 6172200, 9000000, 228600,
            [(str(d.get("foot", "")), 10, False, GRAY, L)])

# ---------- 감사 페이지 (변주 A) ----------
def thanks(slide, d):
    title = str(d.get("title", "감사합니다"))
    tsize = fit_font(title, 10500000, 1, 54, 30)
    textbox(slide, MX, 1828800, 8000000, 274320,
            [(str(d.get("eyebrow", "THANK YOU")), 11, True, BLUE, L)], anchor=MSO_ANCHOR.MIDDLE)
    textbox(slide, MX, 2240280, 10500000, 1100000, [(title, tsize, True, INK, L)])
    box(slide, MX, 3886200, 640080, 64008, fill=BLUE)
    if d.get("subtitle"):
        textbox(slide, MX, 4160520, 9500000, 365760,
                [(str(d["subtitle"]), 18, False, SLATE, L)])
    if d.get("foot"):
        textbox(slide, MX, 6172200, 9000000, 228600,
                [(str(d["foot"]), 10, False, GRAY, L)])

# ---------- 빌드 ----------
def _delete_all_slides(prs):
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(qn('r:id'))
        try: prs.part.drop_rel(rId)
        except Exception: pass
        sldIdLst.remove(sldId)

def build(slides_data, template, out, accent=None):
    """slides_data: [ {cover:true,...}?, {content...}*, {thanks:true,...}? ]
       표지/감사는 페이지 수에서 제외, 본문만 NN/total 번호."""
    configure(accent)
    prs = Presentation(template)
    _delete_all_slides(prs)
    layout = prs.slide_layouts[0]
    content = [d for d in slides_data if not d.get("cover") and not d.get("thanks")]
    total = len(content)
    page = 0
    for d in slides_data:
        slide = prs.slides.add_slide(layout)
        if d.get("cover"):
            cover(slide, d); continue
        if d.get("thanks"):
            thanks(slide, d); continue
        page += 1
        frame(slide, d.get("header", ""), page, total, d.get("caption", ""),
              d.get("title", ""), d.get("subtitle", ""))
        blk = d.get("block")
        if blk and blk.get("type") in BLOCKS:
            BLOCKS[blk["type"]](slide, blk)
    prs.save(out)
    return out
